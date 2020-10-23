import pptx
import pandas as pd

input_pptx_name = 'Example' #name of the input presentation
source_language = ''
target_language = 'fr'

input_pptx = 'input/' + str(input_pptx_name) + '.pptx'

class Pptx_Textract:
    import pptx
    import pandas as pd
    def __init__(self,input_pptx):
        self.input_pptx = input_pptx
        self.prs = pptx.Presentation(input_pptx)
        self.df = pd.DataFrame(columns=['SlideNo','ShapeNo','ParagraphNo','ShapeName','Text','Translation']) #create empty dataframe to store text in
        for slide_nr in range(0, self.prs.slides.__len__()):  # iterate all shapes from all slides
            for shape_nr in range(0, self.prs.slides[slide_nr].shapes.__len__()):
                shape = self.prs.slides[slide_nr].shapes[shape_nr]
                if shape.has_text_frame:  # pull shapes that contain text
                    for i in range (0,shape.text_frame.paragraphs.__len__()):
                        try:#iterate all paragraphs in text frame
                            text = shape.text_frame.paragraphs[i].text  # pull text out of each paragraph
                            self.df = self.df.append(
                                {'SlideNo': slide_nr, 'ShapeNo': shape_nr,  'ParagraphNo': i,'ShapeName': shape.name,
                                 'Text': text,
                                 'Translation': ''}, ignore_index=True)
                        except:
                            pass

        for i in range(0, self.df['Text'].__len__()):  # delete shape from list if empty text contained, else put in <br> identifiers
            if self.df['Text'][i] == '':
                self.df = self.df.drop([i])
            else:
                self.df['Text'][i] = str(self.df['Text'][i]).replace("\n", "<br>")
        self.df = self.df.reset_index(drop=True)

    def to_excel(self,filename):
        self.filename = filename
        timestamp = str(pd.datetime.today())[:16].replace(':', '').replace(' ', '_')
        self.df.to_excel('output/'+ str(self.filename) + '_Textract_' + str(timestamp) +'.xlsx', index=False)  # export dataframe with original text

class GoTra:#google translate api call
    import pandas as pd
    def __init__(self,dataframe,target_language = 'en',source_language = ''):
        self.df = dataframe
        self.source_language = source_language
        self.target_language = target_language
        self.Valid_Lan ={'','de','fr','es','en','nl','da','pl','it','ru','sv'}
        if self.target_language not in self.Valid_Lan or self.source_language not in self.Valid_Lan:
            raise ValueError("Language must be one of %r." % self.Valid_Lan)
    def translate(self):
        import os
        from google.cloud import translate_v2
        os.environ['GOOGLE_APPLICATION_CREDENTIALS'] = "Google_API_Key.json"
        client = translate_v2.Client()
        for i in self.df.index:
            self.df.Translation[i] = str(client.translate(str(self.df.Text[i]),target_language= self.target_language, source_language= self.source_language)['translatedText']).replace('&amp;','&').replace('&quot;','"').replace('&#39;',"'")
    def to_excel(self,filename):
        timestamp = str(pd.datetime.today())[:16].replace(':', '').replace(' ', '_')
        self.filename = filename
        self.df.to_excel('output/'+ str(filename) + '_Translation_' + str(target_language).upper() +'_' + str(timestamp) +'.xlsx', index=False)  # export dataframe with original text

class TransToPptx:
    import pptx
    import pandas as pd
    def import_trans(self, powerpoint, translation):
        prs = pptx.Presentation(powerpoint)
        for i in range(0, translation.__len__()):
                shape = prs.slides[translation['SlideNo'][i]].shapes[translation['ShapeNo'][i]]
                text_frame_paragraph = shape.text_frame.paragraphs[translation['ParagraphNo'][i]]
                trans_text = str(translation['Translation'][i]).replace('<br> ', '\n').replace('<br><br>', '\n\n')
                try:
                    for r in range(0, text_frame_paragraph.runs.__len__()):
                        try:
                            if r > 0:
                                text_frame_paragraph.runs[r].text = ''
                            else:
                                text_frame_paragraph.runs[r].text = trans_text
                            print('working on line ', i, ' run ',r)
                        except:
                            print('Run Error in line', i)
                            pass
                except:
                    print('Error on line ', i, ': --> ', str(text_frame_paragraph.text))
                    pass

        prs.save('output/Translation_' + input_pptx_name + str(pd.datetime.today())[:16].replace(':', '').replace(' ','_')+'.pptx')


pp = Pptx_Textract(input_pptx)
pp.to_excel(input_pptx_name)

translator = GoTra(pp.df,target_language,source_language)
translator.translate()
translator.to_excel(input_pptx_name)

imp = TransToPptx()
imp.import_trans(input_pptx,translator.df)